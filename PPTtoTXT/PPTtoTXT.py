import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation # pip install python-pptx

LANGS = {
    "tr": {
        "lang_label": "Dil seçin:",
        "ppt_btn": "PowerPoint (.pptx) seç",
        "save_btn": "Metin dosyasını kaydet",
        "no_file": "Dosya seçilmedi.",
        "cancel_save": "Kaydetme iptal edildi.",
        "success": "Metinler '{}' dosyasına kaydedildi.",
        "error": "Hata: {}"
    },
    "en": {
        "lang_label": "Select language:",
        "ppt_btn": "Select PowerPoint (.pptx)",
        "save_btn": "Save text file",
        "no_file": "No file selected.",
        "cancel_save": "Saving cancelled.",
        "success": "Texts saved to '{}'.",
        "error": "Error: {}"
    }
}

class App(tk.Tk):
    def __init__(self):
        super().__init__()
        self.title("PPT Text Extractor")
        self.geometry("500x180")
        self.lang_code = "tr"
        self.strings = LANGS[self.lang_code]
        self.pptx_path = None
        self.txt_path = None
        self.create_widgets()

    def create_widgets(self):
        tk.Label(self, text=LANGS["tr"]["lang_label"]).pack(pady=5)

        self.lang_var = tk.StringVar(value="tr")
        lang_menu = tk.OptionMenu(self, self.lang_var, "tr", "en", command=self.change_lang)
        lang_menu.pack()

        self.btn_select_ppt = tk.Button(self, text=self.strings["ppt_btn"], command=self.select_ppt)
        self.btn_select_ppt.pack(pady=10)

        self.btn_save_txt = tk.Button(self, text=self.strings["save_btn"], command=self.save_txt, state="disabled")
        self.btn_save_txt.pack()

        self.status_label = tk.Label(self, text="", fg="green")
        self.status_label.pack(pady=10)

    def change_lang(self, choice):
        self.lang_code = choice
        self.strings = LANGS[self.lang_code]
        self.btn_select_ppt.config(text=self.strings["ppt_btn"])
        self.btn_save_txt.config(text=self.strings["save_btn"])
        self.status_label.config(text="")

    def select_ppt(self):
        path = filedialog.askopenfilename(
            title=self.strings["ppt_btn"],
            filetypes=[(self.strings["ppt_btn"], "*.pptx")]
        )
        if path:
            self.pptx_path = path
            self.status_label.config(text=f"Seçildi: {path}", fg="blue")
            self.btn_save_txt.config(state="normal")
        else:
            self.status_label.config(text=self.strings["no_file"], fg="red")

    def save_txt(self):
        if not self.pptx_path:
            self.status_label.config(text=self.strings["no_file"], fg="red")
            return
        path = filedialog.asksaveasfilename(
            title=self.strings["save_btn"],
            defaultextension=".txt",
            filetypes=[(self.strings["save_btn"], "*.txt")]
        )
        if not path:
            self.status_label.config(text=self.strings["cancel_save"], fg="red")
            return
        try:
            prs = Presentation(self.pptx_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        txt = shape.text_frame.text.strip()
                        if txt:
                            texts.append(txt)
            with open(path, "w", encoding="utf-8") as f:
                f.write("\n\n".join(texts))
            self.status_label.config(text=self.strings["success"].format(path), fg="green")
            messagebox.showinfo(title="Success", message=self.strings["success"].format(path))
        except Exception as e:
            msg = self.strings["error"].format(str(e))
            self.status_label.config(text=msg, fg="red")
            messagebox.showerror(title="Error", message=msg)

if __name__ == "__main__":
    App().mainloop()
